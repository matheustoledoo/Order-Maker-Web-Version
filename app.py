import os
import uuid
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from fpdf import FPDF
import subprocess


app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = os.path.abspath("files")  # Caminho absoluto
app.config["ALLOWED_EXTENSIONS"] = {"pptx"}

# Funções auxiliares
def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in app.config["ALLOWED_EXTENSIONS"]

def aplicar_formatacao(paragraph, fonte="Codec Pro", tamanho=25, cor=(0, 0, 0)):
    if not hasattr(paragraph, "runs"):
        return
    for run in paragraph.runs:
        run.font.name = fonte
        run.font.size = Pt(tamanho)
        run.font.color.rgb = RGBColor(*cor)

def substituir_valores_marcadores(slide, marcador, valor):
    for shape in slide.shapes:
        if not shape.has_text_frame or not isinstance(marcador, str):
            continue
        # Ajustar o tamanho da caixa de texto
        if shape.width < Inches(12):  # Verifica se a largura é menor que o necessário
            shape.width = Inches(12)  # Define uma largura suficiente
        if shape.height < Inches(1):  # Verifica se a altura é menor que o necessário
            shape.height = Inches(1.5)  # Define uma altura suficiente

        for paragraph in shape.text_frame.paragraphs:
            if marcador in paragraph.text:
                paragraph.text = paragraph.text.replace(marcador, valor)
                aplicar_formatacao(paragraph)

def adicionar_lista_incremental(slide, marcador, lista):
    for shape in slide.shapes:
        if not shape.has_text_frame or not isinstance(marcador, str):
            continue
        for paragraph in shape.text_frame.paragraphs:
            if marcador in paragraph.text:
                paragraph.text = marcador
                aplicar_formatacao(paragraph)
                for item in lista:
                    novo_paragraph = shape.text_frame.add_paragraph()
                    novo_paragraph.text = item
                    aplicar_formatacao(novo_paragraph)

def adicionar_equipamentos(slide, lista_equipamentos):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if ":" in paragraph.text:
                paragraph.text += ":"
                aplicar_formatacao(paragraph)
                for equipamento in lista_equipamentos:
                    novo_paragraph = shape.text_frame.add_paragraph()
                    novo_paragraph.text = equipamento
                    aplicar_formatacao(novo_paragraph)
                return

def adicionar_objetos_dinamicos(slide, lista_objetos):
    left = Inches(6)
    top = Inches(3.2)
    width = Inches(5.5)
    height = Inches(0.5)
    espacamento_vertical = Inches(0.9)  # Espaçamento entre linhas

    for obj in lista_objetos:
        if obj.strip():  # Ignora itens vazios
            # Adiciona uma única caixa de texto
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Adiciona o texto respeitando apenas quebras de linha explícitas
            for linha in obj.split("\n"):  # Divide apenas no \n (Enter)
                paragraph = text_frame.add_paragraph()
                paragraph.text = linha.strip()
                aplicar_formatacao(paragraph)

            # Ajusta a posição vertical para o próximo item da lista
            top += espacamento_vertical


def adicionar_escopo_dinamicos(slide, lista_escopo):
    left = Inches(7.1)
    top = Inches(2.6)
    width = Inches(6.5)
    height = Inches(0.5)
    espacamento_vertical = Inches(0.9)  # Espaçamento entre linhas

    for escopo in lista_escopo:
        if escopo.strip():  # Ignora itens vazios
            # Adiciona uma única caixa de texto
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Adiciona o texto respeitando apenas quebras de linha explícitas
            for linha in escopo.split("\n"):  # Divide apenas no \n (Enter)
                paragraph = text_frame.add_paragraph()
                paragraph.text = linha.strip()
                aplicar_formatacao(paragraph)

            # Ajusta a posição vertical para o próximo item da lista
            top += espacamento_vertical



def atualizar_prazo(slide, marcador, valor):
    """
    Substitui o marcador `+` pelo valor fornecido no slide 10,
    mantendo simplicidade e consistência.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:  # Certifica-se de que o shape possui texto
            for paragraph in shape.text_frame.paragraphs:
                if marcador in paragraph.text:  # Verifica se o marcador está no texto
                    paragraph.text = paragraph.text.replace(marcador, valor)  # Substitui o marcador pelo valor
                    aplicar_formatacao(paragraph)  # Aplica a formatação ao texto atualizado


def convert_to_pdf(pptx_path):
    """
    Converte o arquivo PPTX para PDF usando fpdf.
    """
    try:
        prs = Presentation(pptx_path)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        for slide in prs.slides:
            pdf.add_page()
            pdf.set_font("Arial", size=12)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    # Substituir caracteres especiais que causam problemas
                    text = text.replace("–", "-").replace("‘", "'").replace("’", "'").replace("“", '"').replace("”",
                                                                                                                '"')
                    pdf.multi_cell(0, 10, text)

        pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
        pdf.output(pdf_path)
        return pdf_path
    except Exception as e:
        raise Exception(f"Erro ao criar PDF: {e}")

contador = 1

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        try:
            arquivo = request.form["arquivo"]
            caminho_arquivo = os.path.join(app.config["UPLOAD_FOLDER"], arquivo)
            prs = Presentation(caminho_arquivo)

            nome_cliente = request.form.get("nome_cliente", "")
            valor_servico = request.form.get("valor_servico", "")
            valor_mobilizacao = request.form.get("valor_mobilizacao", "")
            prazo = request.form.get("prazo", "")
            objetos = request.form.get("objetos", "").splitlines()
            escopo = request.form.get("escopo", "").splitlines()
            campo = request.form.get("campo", "").splitlines()
            processamento = request.form.get("processamento", "").splitlines()
            equipamentos = request.form.get("equipamentos", "").splitlines()
            texto_slide11 = request.form.get("texto_slide11", "")
            action = request.form.get("action")

            substituir_valores_marcadores(prs.slides[1], "{", nome_cliente)
            substituir_valores_marcadores(prs.slides[1], "}", str(contador))
            substituir_valores_marcadores(prs.slides[10], "{", valor_servico)
            substituir_valores_marcadores(prs.slides[10], "}", valor_mobilizacao)
            adicionar_lista_incremental(prs.slides[7], "Campo", campo)
            adicionar_lista_incremental(prs.slides[7], "Processamento", processamento)
            adicionar_equipamentos(prs.slides[8], equipamentos)
            adicionar_objetos_dinamicos(prs.slides[2], objetos)
            adicionar_escopo_dinamicos(prs.slides[3], escopo)
            atualizar_prazo(prs.slides[10], "+", prazo)



            if texto_slide11.strip():
                for shape in prs.slides[11].shapes:
                    if shape.has_text_frame:
                        shape.text_frame.clear()
                        for linha in texto_slide11.split("\n"):
                            paragraph = shape.text_frame.add_paragraph()
                            paragraph.text = linha
                            aplicar_formatacao(paragraph)

            # Salvar arquivo temporário PPTX
            output_path = os.path.abspath(os.path.join(app.config["UPLOAD_FOLDER"], f"editado_{uuid.uuid4().hex}.pptx"))
            prs.save(output_path)
            contador + 1

            # Converter para PDF se necessário
            if action == "pdf":
                try:
                    output_path = convert_to_pdf(output_path)
                except Exception as e:
                    return f"Erro ao salvar como PDF: {e}"

            return send_file(output_path, as_attachment=True, download_name=os.path.basename(output_path), mimetype="application/pdf" if action == "pdf" else "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as e:
            return f"Erro no processamento: {e}"
        
        

    arquivos = [f for f in os.listdir(app.config["UPLOAD_FOLDER"]) if allowed_file(f)]
    return render_template("index.html", arquivos=arquivos)

if __name__ == "__main__":
    if not os.path.exists(app.config["UPLOAD_FOLDER"]):
        os.makedirs(app.config["UPLOAD_FOLDER"])
    app.run(debug=True)
